"""
Generate the Jobab pitch deck as PPTX with brand styling.

Run:
    python3 scripts/build_pitch_deck.py

Output:
    docs/business/dist/jobab-pitch-deck.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette
GREEN = RGBColor(0x1F, 0x6E, 0x47)
GREEN_DARK = RGBColor(0x15, 0x4D, 0x31)
CREAM = RGBColor(0xFA, 0xF7, 0xF2)
CREAM_2 = RGBColor(0xF4, 0xEF, 0xE6)
INK = RGBColor(0x2A, 0x27, 0x22)
INK_2 = RGBColor(0x6C, 0x64, 0x5A)
AMBER = RGBColor(0x9A, 0x6B, 0x0E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

REPO_ROOT = Path(__file__).resolve().parent.parent
OUT = REPO_ROOT / "docs" / "business" / "dist" / "jobab-pitch-deck.pptx"
LOGO_PNG = REPO_ROOT / "docs" / "public" / "icons" / "logo-512.png"


def add_bg(slide, color):
    """Fill the whole slide with a background color."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, text, left, top, width, height, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font="Hind Siliguri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for line in str(text).split("\n"):
        if p.text == "" and not p.runs:
            run = p.add_run()
        else:
            p = tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, width, height, color=GREEN):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_logo(slide, left, top, height=Inches(0.6)):
    if LOGO_PNG.exists():
        slide.shapes.add_picture(str(LOGO_PNG), left, top, height=height)


def footer(slide, n, total):
    add_text(slide, "Jobab · pre-seed · 2026",
             Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             size=9, color=INK_2)
    add_text(slide, f"{n} / {total}",
             Inches(11.83), Inches(7.05), Inches(1), Inches(0.3),
             size=9, color=INK_2, align=PP_ALIGN.RIGHT)


# ── build ────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

TOTAL = 12

# ── Slide 1 — Title ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0), Inches(0), Inches(13.33), Inches(0.25), GREEN)
add_logo(s, Inches(5.83), Inches(1.6), height=Inches(2.2))
add_text(s, "Jobab", Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.9),
         size=64, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
add_text(s, "AI shop assistant for Bangladeshi merchants",
         Inches(0.5), Inches(4.95), Inches(12.33), Inches(0.6),
         size=22, color=INK_2, align=PP_ALIGN.CENTER)
add_text(s, "[Your name] · Founder · hello@jobab.com",
         Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
         size=14, color=INK_2, align=PP_ALIGN.CENTER)
footer(s, 1, TOTAL)

# ── Slide 2 — Problem ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "The problem", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)
add_text(s, "60% of Bangladeshi small retail happens on Facebook DMs.",
         Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.8),
         size=26, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
add_text(s, "The average shop owner:",
         Inches(0.5), Inches(3.2), Inches(12.33), Inches(0.5),
         size=18, color=INK_2, align=PP_ALIGN.CENTER)
add_text(s, "200+ customer DMs every day",
         Inches(0.5), Inches(3.9), Inches(12.33), Inches(0.5),
         size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "6+ hours/day replying manually",
         Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.5),
         size=24, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "৳15,000+/month in missed sales while sleeping",
         Inches(0.5), Inches(5.1), Inches(12.33), Inches(0.5),
         size=24, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
footer(s, 2, TOTAL)

# ── Slide 3 — Solution ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "The solution", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)
add_text(s, "An AI that answers customer DMs in Bangla, takes orders,",
         Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.7),
         size=24, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "sends bKash links, and hands off to a human when needed.",
         Inches(0.5), Inches(2.5), Inches(12.33), Inches(0.7),
         size=24, color=INK, align=PP_ALIGN.CENTER)
add_accent_bar(s, Inches(1.5), Inches(4.0), Inches(10.33), Inches(2.7), CREAM_2)
add_text(s, "[ screenshot of the Jobab inbox ]",
         Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.5),
         size=14, color=INK_2, align=PP_ALIGN.CENTER)
footer(s, 3, TOTAL)

# ── Slide 4 — How it works (demo) ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "How it works", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

for i, (n, title, body) in enumerate([
    ("1", "Connect your Facebook page", "One click. OAuth. Done."),
    ("2", "AI runs the conversation",
     "Replies in Bangla, recognises products from photos,\ntakes orders, generates bKash payment links."),
    ("3", "You watch live and step in when you want",
     "Take over any thread. Hand back when done."),
]):
    top = Inches(1.8 + i * 1.7)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), top, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    add_text(s, n, Inches(1.0), top + Inches(0.13), Inches(0.9), Inches(0.7),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(2.3), top, Inches(10.5), Inches(0.5),
             size=22, bold=True, color=INK)
    add_text(s, body, Inches(2.3), top + Inches(0.55), Inches(10.5), Inches(1.0),
             size=16, color=INK_2)

add_text(s, "→ live demo at this point",
         Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.3),
         size=12, color=AMBER, align=PP_ALIGN.CENTER)
footer(s, 4, TOTAL)

# ── Slide 5 — Why now ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "Why now", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

for i, (h, body) in enumerate([
    ("LLM cost dropped 10× in 18 months",
     "Llama 3.3 on Groq is fast + cheap enough to make this economical."),
    ("Bangla AI got good in 2025",
     "Llama 4 + Jina v3 finally handle Bangla and Banglish naturally."),
    ("F-commerce is the dominant retail channel",
     "Lazychat raised $1M in 2024 — market validated, still wide open."),
]):
    top = Inches(2.0 + i * 1.5)
    add_accent_bar(s, Inches(1.0), top, Inches(0.08), Inches(1.0), AMBER)
    add_text(s, h, Inches(1.4), top, Inches(11), Inches(0.5),
             size=22, bold=True, color=GREEN_DARK)
    add_text(s, body, Inches(1.4), top + Inches(0.55), Inches(11), Inches(0.5),
             size=16, color=INK_2)
footer(s, 5, TOTAL)

# ── Slide 6 — Market size ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "Market size", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

tiers = [
    ("TAM", "$180M ARR", "500K active F-commerce shops in Bangladesh × $30 ARPU", GREEN),
    ("SAM", "$30M ARR", "50K shops with 100+ DMs/day × $50 ARPU", GREEN_DARK),
    ("SOM (3 years)", "$2.4M ARR", "5K paying shops × $40 ARPU", AMBER),
]
for i, (label, big, sub, color) in enumerate(tiers):
    top = Inches(2.0 + i * 1.5)
    add_accent_bar(s, Inches(1.0), top, Inches(0.12), Inches(1.1), color)
    add_text(s, label, Inches(1.4), top + Inches(0.05), Inches(2.5), Inches(0.5),
             size=16, color=INK_2)
    add_text(s, big, Inches(1.4), top + Inches(0.4), Inches(4), Inches(0.7),
             size=30, bold=True, color=color)
    add_text(s, sub, Inches(5.5), top + Inches(0.4), Inches(7.5), Inches(0.7),
             size=15, color=INK_2)
footer(s, 6, TOTAL)

# ── Slide 7 — Business model ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "Business model", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

tiers = [
    ("Starter", "৳999/mo", "Up to 500 conversations"),
    ("Growth", "৳2,999/mo", "Up to 2,500 + WhatsApp templates"),
    ("Pro", "৳7,999/mo", "Unlimited · multi-page · team seats"),
]
card_w = Inches(3.8)
gap = Inches(0.3)
start_left = (Inches(13.33) - card_w * 3 - gap * 2) / Emu(1)

for i, (name, price, sub) in enumerate(tiers):
    left = Emu(int(start_left + (card_w + gap) / Emu(1) * i))
    top = Inches(2.0)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, Inches(3.0))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE if i != 1 else GREEN
    card.line.color.rgb = GREEN
    card.line.width = Pt(1)
    text_color = INK if i != 1 else WHITE
    sub_color = INK_2 if i != 1 else CREAM_2
    add_text(s, name, left, top + Inches(0.35), card_w, Inches(0.6),
             size=18, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    add_text(s, price, left, top + Inches(1.0), card_w, Inches(0.9),
             size=32, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    add_text(s, sub, left, top + Inches(2.1), card_w, Inches(0.7),
             size=13, color=sub_color, align=PP_ALIGN.CENTER)

add_text(s, "Unit economics  ·  COGS per conversation: $0.003  ·  Gross margin on Growth: ~85%",
         Inches(0.5), Inches(5.6), Inches(12.33), Inches(0.5),
         size=15, color=INK_2, align=PP_ALIGN.CENTER)
footer(s, 7, TOTAL)

# ── Slide 8 — Competition ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "Competition", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

# 2x2 matrix
gx = Inches(3.5)
gy = Inches(2.0)
gw = Inches(7.5)
gh = Inches(4.0)
cell_w = gw / Emu(1) / 2
cell_h = gh / Emu(1) / 2

# Axes labels
add_text(s, "AI depth →", Inches(3.5), Inches(1.6), Inches(7.5), Inches(0.3),
         size=12, color=INK_2, align=PP_ALIGN.CENTER)
add_text(s, "CRM\n↑", Inches(2.5), Inches(3.5), Inches(1.0), Inches(0.8),
         size=12, color=INK_2, align=PP_ALIGN.CENTER)

quads = [
    ("Lazychat", "BD · $1M raised", 0, 0, WHITE),
    ("JOBAB", "(we are here)", 1, 0, GREEN),
    ("Manychat", "Global · generic", 0, 1, WHITE),
    ("AI chatbots", "General-purpose", 1, 1, WHITE),
]
for name, sub, col, row, fill in quads:
    left = Emu(int(gx / Emu(1) + col * cell_w))
    top = Emu(int(gy / Emu(1) + row * cell_h))
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(int(cell_w)), Emu(int(cell_h)))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = GREEN_DARK
    box.line.width = Pt(0.75)
    text_color = WHITE if name == "JOBAB" else INK
    sub_color = CREAM_2 if name == "JOBAB" else INK_2
    add_text(s, name, left, top + Inches(0.4), Emu(int(cell_w)), Inches(0.6),
             size=20, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    add_text(s, sub, left, top + Inches(1.1), Emu(int(cell_w)), Inches(0.5),
             size=13, color=sub_color, align=PP_ALIGN.CENTER)

add_text(s, "We win on Bangla AI depth + photo→product matching + cost telemetry.",
         Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5),
         size=14, color=INK_2, align=PP_ALIGN.CENTER)
footer(s, 8, TOTAL)

# ── Slide 9 — Traction ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "Traction", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

add_text(s, "What's real today", Inches(0.85), Inches(1.6), Inches(6), Inches(0.5),
         size=18, bold=True, color=GREEN_DARK)
shipped = [
    "Full product: inbox, orders, catalog, analytics",
    "Facebook OAuth onboarding",
    "24-hour-window compliance (Meta 2026 rules)",
    "Vision-based product matching",
    "bKash payment link generation",
    "Open-source codebase + docs site",
]
for i, item in enumerate(shipped):
    add_text(s, f"✓  {item}",
             Inches(0.85), Inches(2.2 + i * 0.45), Inches(6), Inches(0.4),
             size=14, color=INK)

add_text(s, "Early signals", Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.5),
         size=18, bold=True, color=GREEN_DARK)
signals = [
    "[N] merchant interviews completed",
    "[N] pilot merchants committed",
    "Meta Business Verification: [date]",
    "App Review: [submitted / pending / approved]",
    "Letters of intent: [merchant 1], [merchant 2]",
]
for i, item in enumerate(signals):
    add_text(s, f"•  {item}",
             Inches(7.0), Inches(2.2 + i * 0.5), Inches(5.8), Inches(0.4),
             size=14, color=INK_2)
footer(s, 9, TOTAL)

# ── Slide 10 — Team ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "Team", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

photo_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.0), Inches(2.5), Inches(2.5))
photo_box.fill.solid()
photo_box.fill.fore_color.rgb = CREAM_2
photo_box.line.color.rgb = GREEN
photo_box.line.width = Pt(2)
add_text(s, "[photo]", Inches(1.5), Inches(3.0), Inches(2.5), Inches(0.5),
         size=12, color=INK_2, align=PP_ALIGN.CENTER)

add_text(s, "[Your name]", Inches(4.5), Inches(2.0), Inches(8), Inches(0.6),
         size=24, bold=True, color=GREEN_DARK)
add_text(s, "Founder", Inches(4.5), Inches(2.6), Inches(8), Inches(0.4),
         size=14, color=AMBER)
bio = (
    "Built Jobab end-to-end (NestJS, Next.js, AI/ML).\n"
    "[Your background — years in tech, prior work, education].\n"
    "Speaks Bangla + English natively."
)
add_text(s, bio, Inches(4.5), Inches(3.1), Inches(8), Inches(1.6),
         size=14, color=INK_2)

add_text(s, "First hires post-funding (12 months)",
         Inches(0.85), Inches(5.3), Inches(12), Inches(0.4),
         size=14, bold=True, color=GREEN_DARK)
add_text(s, "•  Senior full-stack engineer (Dhaka)",
         Inches(1.2), Inches(5.75), Inches(12), Inches(0.35), size=13, color=INK_2)
add_text(s, "•  Merchant success / operations (Dhaka)",
         Inches(1.2), Inches(6.1), Inches(12), Inches(0.35), size=13, color=INK_2)
add_text(s, "•  Fractional CFO (post-MRR milestone)",
         Inches(1.2), Inches(6.45), Inches(12), Inches(0.35), size=13, color=INK_2)
footer(s, 10, TOTAL)

# ── Slide 11 — The ask ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bar(s, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), GREEN)
add_text(s, "The ask", Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
         size=32, bold=True, color=INK)

add_text(s, "$400K seed", Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.2),
         size=72, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
add_text(s, "18-month runway",
         Inches(0.5), Inches(3.0), Inches(12.33), Inches(0.5),
         size=20, color=INK_2, align=PP_ALIGN.CENTER)

add_text(s, "Milestones", Inches(0.85), Inches(4.0), Inches(6), Inches(0.4),
         size=16, bold=True, color=GREEN_DARK)
for i, m in enumerate([
    "100 paying merchants by month 12",
    "$50K MRR by month 18",
    "App Review approved + WhatsApp Cloud live by month 3",
    "Series A ready by month 18",
]):
    add_text(s, f"•  {m}", Inches(1.2), Inches(4.45 + i * 0.4), Inches(6), Inches(0.35),
             size=13, color=INK)

add_text(s, "Use of funds", Inches(7.0), Inches(4.0), Inches(5.8), Inches(0.4),
         size=16, bold=True, color=GREEN_DARK)
for i, (label, pct) in enumerate([
    ("Engineering (1 senior eng)", "40%"),
    ("Merchant operations", "25%"),
    ("Marketing + outreach", "15%"),
    ("Infrastructure", "15%"),
    ("Buffer", "5%"),
]):
    add_text(s, f"•  {label}", Inches(7.3), Inches(4.45 + i * 0.4), Inches(4.5), Inches(0.35),
             size=13, color=INK)
    add_text(s, pct, Inches(11.8), Inches(4.45 + i * 0.4), Inches(0.9), Inches(0.35),
             size=13, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
footer(s, 11, TOTAL)

# ── Slide 12 — Thank you ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN_DARK)
add_logo(s, Inches(6.16), Inches(0.7), height=Inches(1.6))
add_text(s, "Let's give 50,000 Bangladeshi",
         Inches(0.5), Inches(2.6), Inches(12.33), Inches(0.8),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "shop owners their evenings back.",
         Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.8),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "[Your name]  ·  hello@jobab.com  ·  +880 1XXX-XXXXXX",
         Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.5),
         size=18, color=CREAM_2, align=PP_ALIGN.CENTER)
add_text(s, "Live docs: siraajul.github.io/jobab",
         Inches(0.5), Inches(5.7), Inches(12.33), Inches(0.5),
         size=14, color=CREAM, align=PP_ALIGN.CENTER)
add_text(s, "Repo: github.com/siraajul/jobab",
         Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.5),
         size=14, color=CREAM, align=PP_ALIGN.CENTER)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"✓ wrote {OUT.relative_to(REPO_ROOT)} ({OUT.stat().st_size // 1024} KB)")
